"""
file_extractor_service.py – Extract text content from various file formats.

Supported formats:
  Text:   .txt, .md, .rst, .log, .csv, .tsv, .json, .jsonl, .xml, .yaml, .yml,
          .ini, .cfg, .toml, .html, .htm, .sql, and common code extensions
  PDF:    PyMuPDF (fitz)
  DOCX:   python-docx
  PPTX:   python-pptx
  XLSX:   openpyxl (read-only)
  HTML:   BeautifulSoup
  Image:  Base64 Data URL for Vision API (or skip)
  Audio:  faster-whisper + FFmpeg
  Video:  FFmpeg extract audio → faster-whisper
  ZIP:    Safe extract → recurse into supported files
"""
from __future__ import annotations

import base64
import hashlib
import io
import mimetypes
import shutil
import subprocess
import tempfile
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

from settings import (
    AI_XLSX_MAX_COLS,
    AI_XLSX_MAX_ROWS,
    AI_ZIP_MAX_EXTRACT_MB,
    AI_ZIP_MAX_FILES,
)


# ── Result dataclass ──────────────────────────────────────────────────────────

@dataclass
class ExtractionResult:
    text: Optional[str] = None
    image_data_url: Optional[str] = None   # for Vision mode
    extraction_method: str = "unknown"
    page_count: int = 0
    sheet_count: int = 0
    slide_count: int = 0
    error_code: Optional[str] = None
    error_message: Optional[str] = None
    metadata: dict = field(default_factory=dict)

    @property
    def success(self) -> bool:
        return self.error_code is None

    @property
    def has_content(self) -> bool:
        return bool(self.text) or bool(self.image_data_url)


# ── Dispatcher ────────────────────────────────────────────────────────────────

def extract_file(
    path: Path,
    *,
    vision_enabled: bool = False,
    audio_mode: str = "transcribe_and_send",  # transcribe_and_send | transcript_only | skip
    zip_enabled: bool = False,
    intermediate_dir: Optional[Path] = None,
) -> ExtractionResult:
    """Main dispatcher: detect file type and delegate to the appropriate extractor."""
    if not path.exists():
        return ExtractionResult(error_code="FILE_NOT_FOUND", error_message=f"فایل یافت نشد: {path}")

    ext = path.suffix.lower()

    # ── Text / Code
    if ext in _TEXT_EXTS:
        return _extract_text(path)

    # ── HTML (separate from generic text for tag stripping)
    if ext in (".html", ".htm"):
        return _extract_html(path)

    # ── PDF
    if ext == ".pdf":
        return _extract_pdf(path)

    # ── DOCX
    if ext in (".docx", ".doc"):
        return _extract_docx(path)

    # ── PPTX
    if ext in (".pptx", ".ppt"):
        return _extract_pptx(path)

    # ── XLSX / XLS
    if ext in (".xlsx", ".xls"):
        return _extract_xlsx(path)

    # ── CSV / TSV (special handling for delimiter detection)
    if ext in (".csv", ".tsv"):
        return _extract_csv(path)

    # ── Image
    if ext in (".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif", ".tiff", ".tif"):
        if vision_enabled:
            return _extract_image_vision(path)
        return ExtractionResult(
            error_code="VISION_NOT_SUPPORTED",
            error_message="پردازش تصویر غیرفعال است. Vision را در تنظیمات فعال کنید.",
            extraction_method="skipped",
        )

    # ── Audio / Video
    if ext in (".mp3", ".m4a", ".wav", ".aac", ".flac", ".ogg", ".opus",
               ".mp4", ".mkv", ".mov", ".webm", ".avi", ".flv"):
        if audio_mode == "skip":
            return ExtractionResult(
                error_code="SKIPPED",
                error_message="فایل صوتی/ویدیویی Skip شد (تنظیم کاربر).",
                extraction_method="skipped",
            )
        return _extract_audio_video(path, intermediate_dir=intermediate_dir)

    # ── ZIP
    if ext == ".zip":
        if zip_enabled:
            return _extract_zip(path, intermediate_dir=intermediate_dir)
        return ExtractionResult(
            error_code="SKIPPED",
            error_message="پردازش ZIP غیرفعال است.",
            extraction_method="skipped",
        )

    return ExtractionResult(
        error_code="UNSUPPORTED_FILE_TYPE",
        error_message=f"نوع فایل پشتیبانی نمی‌شود: {ext}",
        extraction_method="unsupported",
    )


# ── Text ──────────────────────────────────────────────────────────────────────

_TEXT_EXTS = frozenset({
    ".txt", ".md", ".rst", ".log", ".json", ".jsonl", ".xml",
    ".yaml", ".yml", ".ini", ".cfg", ".toml", ".sql",
    ".py", ".js", ".ts", ".tsx", ".jsx", ".java", ".cs",
    ".cpp", ".c", ".h", ".hpp", ".go", ".rs", ".php",
    ".rb", ".swift", ".kt", ".dart", ".sh", ".ps1", ".bat",
})


def _detect_encoding(data: bytes) -> str:
    """Detect encoding with charset-normalizer, fallback to utf-8."""
    try:
        from charset_normalizer import from_bytes
        result = from_bytes(data).best()
        return result.encoding if result else "utf-8"
    except ImportError:
        return "utf-8"


def _extract_text(path: Path) -> ExtractionResult:
    try:
        raw = path.read_bytes()
        if _looks_binary(raw[:1024]):
            return ExtractionResult(
                error_code="TEXT_EXTRACTION_FAILED",
                error_message="فایل باینری به اشتباه به عنوان متنی شناسایی شد.",
                extraction_method="direct_text",
            )
        enc = "utf-8"
        try:
            text = raw.decode("utf-8")
        except UnicodeDecodeError:
            enc = _detect_encoding(raw)
            text = raw.decode(enc, errors="replace")
        return ExtractionResult(
            text=text,
            extraction_method="direct_text",
            metadata={"encoding": enc, "size_bytes": len(raw)},
        )
    except PermissionError:
        return ExtractionResult(error_code="PERMISSION_DENIED", error_message="دسترسی به فایل ممنوع است.")
    except OSError as exc:
        return ExtractionResult(error_code="TEXT_EXTRACTION_FAILED", error_message=str(exc))


# ── HTML ─────────────────────────────────────────────────────────────────────

def _extract_html(path: Path) -> ExtractionResult:
    text_result = _extract_text(path)
    if not text_result.success or not text_result.text:
        return text_result

    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(text_result.text, "html.parser")
        for tag in soup(["script", "style", "noscript", "head"]):
            tag.decompose()
        title = soup.title.string.strip() if soup.title and soup.title.string else ""
        lines = []
        if title:
            lines.append(f"# {title}\n")
        lines.append(soup.get_text(separator="\n", strip=True))
        return ExtractionResult(
            text="\n".join(lines),
            extraction_method="beautifulsoup",
            metadata={"title": title},
        )
    except ImportError:
        # Fallback: return raw text
        return ExtractionResult(
            text=text_result.text,
            extraction_method="direct_text_html",
        )
    except Exception as exc:
        return ExtractionResult(
            error_code="TEXT_EXTRACTION_FAILED",
            error_message=f"خطا در پردازش HTML: {exc}",
        )


# ── CSV / TSV ────────────────────────────────────────────────────────────────

def _extract_csv(path: Path) -> ExtractionResult:
    import csv as csv_mod

    text_result = _extract_text(path)
    if not text_result.success or not text_result.text:
        return text_result

    try:
        sample = text_result.text[:8192]
        dialect = csv_mod.Sniffer().sniff(sample, delimiters=",\t;|")
        delimiter = dialect.delimiter
    except csv_mod.Error:
        delimiter = "\t" if path.suffix.lower() == ".tsv" else ","

    try:
        reader = csv_mod.reader(io.StringIO(text_result.text), delimiter=delimiter)
        rows = []
        for i, row in enumerate(reader):
            if i >= AI_XLSX_MAX_ROWS:
                rows.append(f"[... {i} ردیف نمایش داده شد، بقیه برش خورد]")
                break
            rows.append("\t".join(row))
        return ExtractionResult(
            text="\n".join(rows),
            extraction_method="csv_reader",
            metadata={"delimiter": repr(delimiter)},
        )
    except Exception as exc:
        return ExtractionResult(
            error_code="TEXT_EXTRACTION_FAILED",
            error_message=f"خطا در پردازش CSV: {exc}",
        )


# ── PDF ───────────────────────────────────────────────────────────────────────

def _extract_pdf(path: Path) -> ExtractionResult:
    try:
        import fitz  # PyMuPDF
    except ImportError:
        return ExtractionResult(
            error_code="TEXT_EXTRACTION_FAILED",
            error_message="PyMuPDF نصب نشده است. دستور: pip install pymupdf",
        )

    try:
        doc = fitz.open(str(path))
        if doc.is_encrypted:
            doc.close()
            return ExtractionResult(
                error_code="ENCRYPTED_FILE",
                error_message="فایل PDF رمزدار است.",
                extraction_method="pymupdf",
                metadata={"page_count": 0},
            )

        pages: list[str] = []
        for i, page in enumerate(doc):
            text = page.get_text("text")
            pages.append(f"--- Page {i + 1} ---\n{text.strip()}")

        doc.close()
        full_text = "\n\n".join(pages)

        if not full_text.strip():
            return ExtractionResult(
                error_code="TEXT_EXTRACTION_FAILED",
                error_message="PDF متن قابل‌استخراج ندارد. ممکن است اسکن‌شده باشد.",
                extraction_method="pymupdf",
                metadata={"page_count": len(pages)},
            )

        return ExtractionResult(
            text=full_text,
            extraction_method="pymupdf",
            page_count=len(pages),
            metadata={"page_count": len(pages)},
        )
    except Exception as exc:
        return ExtractionResult(
            error_code="TEXT_EXTRACTION_FAILED",
            error_message=f"خطا در پردازش PDF: {exc}",
            extraction_method="pymupdf",
        )


# ── DOCX ──────────────────────────────────────────────────────────────────────

def _extract_docx(path: Path) -> ExtractionResult:
    try:
        from docx import Document
        from docx.oxml.ns import qn
    except ImportError:
        return ExtractionResult(
            error_code="TEXT_EXTRACTION_FAILED",
            error_message="python-docx نصب نشده است. دستور: pip install python-docx",
        )

    try:
        doc = Document(str(path))
        lines: list[str] = []

        for para in doc.paragraphs:
            style = para.style.name if para.style else ""
            text = para.text.strip()
            if not text:
                continue
            if style.startswith("Heading"):
                level = style.replace("Heading", "").strip()
                prefix = "#" * (int(level) if level.isdigit() else 1)
                lines.append(f"{prefix} {text}")
            else:
                lines.append(text)

        for table in doc.tables:
            lines.append("")
            for row in table.rows:
                lines.append("\t".join(cell.text.strip() for cell in row.cells))
            lines.append("")

        return ExtractionResult(
            text="\n".join(lines),
            extraction_method="python_docx",
            metadata={"paragraphs": len(doc.paragraphs), "tables": len(doc.tables)},
        )
    except Exception as exc:
        return ExtractionResult(
            error_code="TEXT_EXTRACTION_FAILED",
            error_message=f"خطا در پردازش DOCX: {exc}",
        )


# ── PPTX ──────────────────────────────────────────────────────────────────────

def _extract_pptx(path: Path) -> ExtractionResult:
    try:
        from pptx import Presentation
    except ImportError:
        return ExtractionResult(
            error_code="TEXT_EXTRACTION_FAILED",
            error_message="python-pptx نصب نشده است. دستور: pip install python-pptx",
        )

    try:
        prs = Presentation(str(path))
        slides_text: list[str] = []

        for i, slide in enumerate(prs.slides):
            parts: list[str] = [f"--- Slide {i + 1} ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
            # Speaker notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    parts.append(f"[Notes] {notes_text}")
            slides_text.append("\n".join(parts))

        return ExtractionResult(
            text="\n\n".join(slides_text),
            extraction_method="python_pptx",
            slide_count=len(prs.slides),
            metadata={"slide_count": len(prs.slides)},
        )
    except Exception as exc:
        return ExtractionResult(
            error_code="TEXT_EXTRACTION_FAILED",
            error_message=f"خطا در پردازش PPTX: {exc}",
        )


# ── XLSX ──────────────────────────────────────────────────────────────────────

def _extract_xlsx(path: Path) -> ExtractionResult:
    try:
        import openpyxl
    except ImportError:
        return ExtractionResult(
            error_code="TEXT_EXTRACTION_FAILED",
            error_message="openpyxl نصب نشده است. دستور: pip install openpyxl",
        )

    try:
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        sheets_text: list[str] = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows_text: list[str] = [f"--- Sheet: {sheet_name} ---"]
            row_count = 0
            for row in ws.iter_rows(max_row=AI_XLSX_MAX_ROWS, max_col=AI_XLSX_MAX_COLS):
                cells = [str(cell.value) if cell.value is not None else "" for cell in row]
                rows_text.append("\t".join(cells))
                row_count += 1
            if row_count >= AI_XLSX_MAX_ROWS:
                rows_text.append(f"[... بیش از {AI_XLSX_MAX_ROWS} ردیف، برش خورد]")
            sheets_text.append("\n".join(rows_text))

        wb.close()
        return ExtractionResult(
            text="\n\n".join(sheets_text),
            extraction_method="openpyxl",
            sheet_count=len(wb.sheetnames),
            metadata={"sheet_count": len(wb.sheetnames)},
        )
    except Exception as exc:
        return ExtractionResult(
            error_code="TEXT_EXTRACTION_FAILED",
            error_message=f"خطا در پردازش XLSX: {exc}",
        )


# ── Image (Vision) ────────────────────────────────────────────────────────────

def _extract_image_vision(path: Path) -> ExtractionResult:
    try:
        raw = path.read_bytes()
        mime = mimetypes.guess_type(str(path))[0] or "image/jpeg"
        b64 = base64.b64encode(raw).decode("ascii")
        data_url = f"data:{mime};base64,{b64}"
        return ExtractionResult(
            image_data_url=data_url,
            extraction_method="vision_base64",
            metadata={"mime": mime, "size_bytes": len(raw)},
        )
    except OSError as exc:
        return ExtractionResult(
            error_code="TEXT_EXTRACTION_FAILED",
            error_message=f"خطا در خواندن تصویر: {exc}",
        )


# ── Audio / Video ─────────────────────────────────────────────────────────────

def _extract_audio_video(path: Path, intermediate_dir: Optional[Path] = None) -> ExtractionResult:
    """Transcribe audio/video with faster-whisper. Returns transcript text."""
    # Check for cached transcript
    if intermediate_dir:
        cache_path = intermediate_dir / f"{path.stem}__transcript.txt"
        if cache_path.exists() and cache_path.stat().st_size > 0:
            try:
                text = cache_path.read_text(encoding="utf-8")
                return ExtractionResult(
                    text=text,
                    extraction_method="whisper_cached",
                    metadata={"transcript_cache": str(cache_path)},
                )
            except OSError:
                pass

    # Extract audio if video
    audio_path = path
    tmp_dir: Optional[tempfile.TemporaryDirectory] = None
    is_video = path.suffix.lower() in (".mp4", ".mkv", ".mov", ".webm", ".avi", ".flv")

    if is_video:
        try:
            from settings import get_ffmpeg_executable
            ffmpeg_exe = get_ffmpeg_executable()
            tmp_dir = tempfile.TemporaryDirectory()
            audio_path = Path(tmp_dir.name) / f"{path.stem}.wav"
            subprocess.run(
                [ffmpeg_exe, "-y", "-i", str(path), "-vn", "-ac", "1",
                 "-ar", "16000", str(audio_path)],
                check=True,
                capture_output=True,
                timeout=600,
            )
        except FileNotFoundError:
            return ExtractionResult(
                error_code="FFMPEG_FAILED",
                error_message="FFmpeg یافت نشد. برای پردازش ویدیو FFmpeg را نصب کنید.",
            )
        except subprocess.CalledProcessError as exc:
            return ExtractionResult(
                error_code="FFMPEG_FAILED",
                error_message=f"خطا در استخراج صدا: {exc.stderr.decode(errors='replace')[:300]}",
            )
        except subprocess.TimeoutExpired:
            return ExtractionResult(error_code="FFMPEG_FAILED", error_message="Timeout در استخراج صدا.")

    # Transcribe
    try:
        from faster_whisper import WhisperModel
        from settings import (
            WHISPER_BEAM_SIZE,
            WHISPER_COMPUTE_TYPE,
            WHISPER_DEVICE,
            WHISPER_LANGUAGE,
            WHISPER_MODEL,
            WHISPER_VAD_FILTER,
        )

        model = WhisperModel(
            WHISPER_MODEL,
            device=WHISPER_DEVICE if WHISPER_DEVICE != "auto" else "cpu",
            compute_type=WHISPER_COMPUTE_TYPE if WHISPER_COMPUTE_TYPE != "auto" else "int8",
        )
        segments, info = model.transcribe(
            str(audio_path),
            beam_size=WHISPER_BEAM_SIZE,
            language=WHISPER_LANGUAGE or None,
            vad_filter=WHISPER_VAD_FILTER,
        )
        lines = [seg.text.strip() for seg in segments]
        text = "\n".join(lines)

        if intermediate_dir:
            try:
                intermediate_dir.mkdir(parents=True, exist_ok=True)
                cache_path = intermediate_dir / f"{path.stem}__transcript.txt"
                cache_path.write_text(text, encoding="utf-8")
            except OSError:
                pass

        return ExtractionResult(
            text=text,
            extraction_method="faster_whisper",
            metadata={
                "language_detected": getattr(info, "language", ""),
                "duration_s": getattr(info, "duration", 0),
            },
        )
    except ImportError:
        return ExtractionResult(
            error_code="WHISPER_FAILED",
            error_message="faster-whisper نصب نشده است.",
        )
    except Exception as exc:
        return ExtractionResult(
            error_code="WHISPER_FAILED",
            error_message=f"خطا در تبدیل صوت: {exc}",
        )
    finally:
        if tmp_dir:
            try:
                tmp_dir.cleanup()
            except Exception:
                pass


# ── ZIP ───────────────────────────────────────────────────────────────────────

def _extract_zip(path: Path, intermediate_dir: Optional[Path] = None) -> ExtractionResult:
    """Safely extract and collect text from supported files inside a ZIP."""
    try:
        with zipfile.ZipFile(str(path), "r") as zf:
            members = zf.namelist()
            if len(members) > AI_ZIP_MAX_FILES:
                return ExtractionResult(
                    error_code="TEXT_EXTRACTION_FAILED",
                    error_message=f"ZIP دارای بیش از {AI_ZIP_MAX_FILES} فایل است.",
                )

            with tempfile.TemporaryDirectory() as tmp:
                tmp_path = Path(tmp)
                # Zip Slip check
                for member in members:
                    member_path = (tmp_path / member).resolve()
                    if not str(member_path).startswith(str(tmp_path.resolve())):
                        return ExtractionResult(
                            error_code="TEXT_EXTRACTION_FAILED",
                            error_message=f"ZIP Slip تشخیص داده شد: {member}",
                        )

                # Check uncompressed size
                total_bytes = sum(info.file_size for info in zf.infolist())
                max_bytes = AI_ZIP_MAX_EXTRACT_MB * 1024 * 1024
                if total_bytes > max_bytes:
                    return ExtractionResult(
                        error_code="TEXT_EXTRACTION_FAILED",
                        error_message=f"حجم Extract شده بیش از {AI_ZIP_MAX_EXTRACT_MB} MB است.",
                    )

                zf.extractall(tmp)
                all_texts: list[str] = []

                for member in sorted(members):
                    member_path = tmp_path / member
                    if not member_path.is_file():
                        continue
                    inner = extract_file(member_path, intermediate_dir=intermediate_dir)
                    if inner.success and inner.text:
                        all_texts.append(f"=== {member} ===\n{inner.text}")

                return ExtractionResult(
                    text="\n\n".join(all_texts),
                    extraction_method="zip_extract",
                    metadata={"member_count": len(members), "extracted_files": len(all_texts)},
                )
    except zipfile.BadZipFile:
        return ExtractionResult(error_code="TEXT_EXTRACTION_FAILED", error_message="فایل ZIP خراب است.")
    except Exception as exc:
        return ExtractionResult(
            error_code="TEXT_EXTRACTION_FAILED",
            error_message=f"خطا در پردازش ZIP: {exc}",
        )


# ── Helpers ───────────────────────────────────────────────────────────────────

def _looks_binary(sample: bytes) -> bool:
    if not sample:
        return False
    non_printable = sum(1 for b in sample if b < 9 or (13 < b < 32) or b == 127)
    return non_printable / len(sample) > 0.30
